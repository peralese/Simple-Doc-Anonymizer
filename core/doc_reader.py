"""Format-aware document reader — returns structured content for the pipeline."""

from pathlib import Path
from typing import Any


def read(path: str) -> dict[str, Any]:
    """
    Read a document and return a format-agnostic payload dict:
      {
        "format": str,
        "content": <format-specific structure>,
        "meta": dict          # human-readable stats for the progress line
      }
    """
    p = Path(path)
    ext = p.suffix.lower()

    if ext in (".txt", ".md"):
        return _read_text(p)
    elif ext == ".docx":
        return _read_docx(p)
    elif ext in (".xlsx", ".xls"):
        return _read_xlsx(p)
    elif ext == ".csv":
        return _read_csv(p)
    elif ext == ".pdf":
        return _read_pdf(p)
    elif ext == ".pptx":
        return _read_pptx(p)
    else:
        raise ValueError(f"Unsupported file format: {ext}")


def _read_text(p: Path) -> dict:
    text = p.read_text(encoding="utf-8-sig")
    return {
        "format": "text",
        "content": text,
        "meta": {"lines": text.count("\n") + 1},
    }


def _read_docx(p: Path) -> dict:
    from docx import Document
    doc = Document(str(p))
    paragraphs = [para.text for para in doc.paragraphs]
    return {
        "format": "docx",
        "content": {"doc": doc, "paragraphs": paragraphs},
        "meta": {"paragraphs": len(paragraphs)},
    }


def _read_xlsx(p: Path) -> dict:
    import openpyxl
    wb = openpyxl.load_workbook(str(p))
    sheets: dict[str, list[list[Any]]] = {}
    total_cells = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows():
            rows.append([(cell.row, cell.column, cell.value) for cell in row])
            total_cells += sum(1 for cell in row if cell.value is not None)
        sheets[sheet_name] = rows
    return {
        "format": "xlsx",
        "content": {"workbook": wb, "sheets": sheets},
        "meta": {"sheets": len(wb.sheetnames), "cells": total_cells},
    }


def _read_csv(p: Path) -> dict:
    import csv
    rows: list[list[str]] = []
    with p.open(newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(row)
    return {
        "format": "csv",
        "content": rows,
        "meta": {"rows": len(rows), "cols": len(rows[0]) if rows else 0},
    }


def _read_pdf(p: Path) -> dict:
    import pdfplumber
    pages_text: list[str] = []
    with pdfplumber.open(str(p)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            pages_text.append(text)
    full_text = "\n".join(pages_text)
    return {
        "format": "pdf",
        "content": full_text,
        "meta": {"pages": len(pages_text)},
    }


def _read_pptx(p: Path) -> dict:
    from pptx import Presentation
    prs = Presentation(str(p))
    slides_info: list[dict] = []
    total_shapes = 0
    for slide_idx, slide in enumerate(prs.slides):
        shapes_info = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                total_shapes += 1
                shapes_info.append({
                    "shape_idx": slide.shapes.index(shape),
                    "shape_name": shape.name,
                    "paragraphs": [para.text for para in shape.text_frame.paragraphs],
                })
        slides_info.append({"slide_idx": slide_idx, "shapes": shapes_info})
    return {
        "format": "pptx",
        "content": {"presentation": prs, "slides": slides_info},
        "meta": {"slides": len(prs.slides), "text_shapes": total_shapes},
    }
