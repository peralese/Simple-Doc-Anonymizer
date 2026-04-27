"""Format-aware document writer — writes redacted content to disk."""

import csv
from pathlib import Path
from typing import Any


def write(
    original_payload: dict[str, Any],
    redacted_content: Any,
    output_path: str,
) -> str:
    """
    Write redacted content to output_path.
    Returns the actual path written (may differ for PDF → txt).
    """
    fmt = original_payload["format"]
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    if fmt == "text":
        return _write_text(redacted_content, out)
    elif fmt == "docx":
        return _write_docx(original_payload, redacted_content, out)
    elif fmt == "xlsx":
        return _write_xlsx(original_payload, redacted_content, out)
    elif fmt == "csv":
        return _write_csv(redacted_content, out)
    elif fmt == "pdf":
        return _write_pdf_as_txt(redacted_content, out)
    elif fmt == "pptx":
        return _write_pptx(original_payload, redacted_content, out)
    else:
        raise ValueError(f"Unsupported format for writing: {fmt}")


def _write_text(text: str, out: Path) -> str:
    out.write_text(text, encoding="utf-8")
    return str(out)


def _write_docx(payload: dict, redacted_paragraphs: list[str], out: Path) -> str:
    doc = payload["content"]["doc"]
    for para, new_text in zip(doc.paragraphs, redacted_paragraphs):
        if para.runs:
            # Preserve formatting of the first run, clear the rest
            para.runs[0].text = new_text
            for run in para.runs[1:]:
                run.text = ""
        else:
            para.clear()
            para.add_run(new_text)
    doc.save(str(out))
    return str(out)


def _write_xlsx(payload: dict, redacted_sheets: dict[str, list[list[Any]]], out: Path) -> str:
    wb = payload["content"]["workbook"]
    for sheet_name, redacted_rows in redacted_sheets.items():
        ws = wb[sheet_name]
        original_rows = payload["content"]["sheets"][sheet_name]
        for orig_row, redacted_row in zip(original_rows, redacted_rows):
            for (r, c, _orig), new_val in zip(orig_row, redacted_row):
                ws.cell(row=r, column=c).value = new_val
    wb.save(str(out))
    return str(out)


def _write_csv(rows: list[list[str]], out: Path) -> str:
    with out.open("w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerows(rows)
    return str(out)


def _write_pdf_as_txt(text: str, out: Path) -> str:
    # PDF write-back is not supported; output as .txt
    txt_out = out.with_suffix(".txt")
    header = (
        "# NOTE: PDF write-back is not supported.\n"
        "# This file contains the extracted and redacted text.\n\n"
    )
    txt_out.write_text(header + text, encoding="utf-8")
    return str(txt_out)


def _write_pptx(payload: dict, redacted_slides: list[dict], out: Path) -> str:
    prs = payload["content"]["presentation"]
    for slide_info, slide in zip(redacted_slides, prs.slides):
        for shape_info in slide_info["shapes"]:
            shape = slide.shapes[shape_info["shape_idx"]]
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for para, new_text in zip(tf.paragraphs, shape_info["paragraphs"]):
                if para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    # Add a run if there are none
                    from pptx.util import Pt  # noqa: F401
                    run = para.add_run()
                    run.text = new_text
    prs.save(str(out))
    return str(out)
