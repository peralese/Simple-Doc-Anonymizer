"""Format-aware document reader for doc-anonymizer-v2.

Returns per-chunk text with location metadata for the detection pipeline.
"""

from pathlib import Path


def read_document(path: str) -> dict:
    """
    Read a document and return:
    {
        "format": str,
        "chunks": [{"text": str, "location": str}],
        "raw": <format-specific object used by doc_writer>,
        "meta": dict
    }
    """
    p = Path(path)
    ext = p.suffix.lower()
    dispatch = {
        ".txt":  _read_text,
        ".md":   _read_text,
        ".docx": _read_docx,
        ".xlsx": _read_xlsx,
        ".xls":  _read_xlsx,
        ".csv":  _read_csv,
        ".pdf":  _read_pdf,
        ".pptx": _read_pptx,
    }
    reader = dispatch.get(ext)
    if reader is None:
        raise ValueError(f"Unsupported file format: {ext!r}")
    return reader(p)


# ── Format readers ────────────────────────────────────────────────────────────

def _read_text(p: Path) -> dict:
    text = p.read_text(encoding="utf-8-sig")
    lines = text.splitlines()
    chunks = [
        {"text": line, "location": f"Line {i + 1}"}
        for i, line in enumerate(lines)
        if line.strip()
    ]
    return {
        "format": "text",
        "chunks": chunks,
        "raw": text,
        "meta": {"lines": len(lines)},
    }


def _read_docx(p: Path) -> dict:
    from docx import Document
    doc = Document(str(p))
    chunks = [
        {"text": para.text, "location": f"Para {i + 1}"}
        for i, para in enumerate(doc.paragraphs)
        if para.text.strip()
    ]
    return {
        "format": "docx",
        "chunks": chunks,
        "raw": doc,
        "meta": {"paragraphs": len(doc.paragraphs)},
    }


def _read_xlsx(p: Path) -> dict:
    import openpyxl
    from openpyxl.utils import get_column_letter

    wb = openpyxl.load_workbook(str(p))
    chunks = []
    total_cells = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None and str(cell.value).strip():
                    col_letter = get_column_letter(cell.column)
                    location = f"{sheet_name}!{col_letter}{cell.row}"
                    chunks.append({"text": str(cell.value), "location": location})
                    total_cells += 1

    return {
        "format": "xlsx",
        "chunks": chunks,
        "raw": wb,
        "meta": {"sheets": len(wb.sheetnames), "cells": total_cells},
    }


def _read_csv(p: Path) -> dict:
    import csv

    chunks = []
    rows_data: list[list[str]] = []

    with p.open(newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        for row_idx, row in enumerate(reader, start=1):
            rows_data.append(row)
            for col_idx, cell in enumerate(row, start=1):
                if cell.strip():
                    chunks.append({
                        "text": cell,
                        "location": f"Row {row_idx}, Col {col_idx}",
                    })

    return {
        "format": "csv",
        "chunks": chunks,
        "raw": rows_data,
        "meta": {"rows": len(rows_data), "cols": len(rows_data[0]) if rows_data else 0},
    }


def _read_pdf(p: Path) -> dict:
    import pdfplumber

    chunks = []
    page_texts: list[str] = []

    with pdfplumber.open(str(p)) as pdf:
        for page_idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            page_texts.append(text)
            for line_idx, line in enumerate(text.splitlines(), start=1):
                if line.strip():
                    chunks.append({
                        "text": line,
                        "location": f"Page {page_idx}, Line {line_idx}",
                    })

    full_text = "\n".join(page_texts)
    return {
        "format": "pdf",
        "chunks": chunks,
        "raw": full_text,   # full extracted text; PDF write-back uses this as .txt
        "meta": {"pages": len(page_texts)},
    }


def _read_pptx(p: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(p))
    chunks = []
    total_shapes = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            total_shapes += 1
            location = f"Slide {slide_idx} / Shape: {shape.name}"
            full_text = shape.text_frame.text
            if full_text.strip():
                chunks.append({"text": full_text, "location": location})

    return {
        "format": "pptx",
        "chunks": chunks,
        "raw": prs,
        "meta": {"slides": len(prs.slides), "text_shapes": total_shapes},
    }
